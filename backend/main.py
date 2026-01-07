from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import requests
import uuid
from io import BytesIO
import math

# =========================
# Utils
# =========================
from utils.logger import get_logger
from utils.chunker import chunk_text
from utils.embedder import embed_chunks
from utils.retriever import retrieve_relevant_chunks
from utils.prompt_builder import build_rag_prompt
from utils.timer import timed

# =========================
# App Initialization
# =========================

app = FastAPI(
    title="Cortex Backend",
    version="1.2.2",
)

logger = get_logger("main")

# =========================
# CORS
# =========================

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =========================
# Ollama Config
# =========================

OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "llama3.1:8b"

# =========================
# In-Memory Ephemeral Store
# =========================
# session_id → {
#   type: "document" | "image",
#   chunks?: list[str],
#   embeddings?: list[list[float]],
#   ocr_text?: str
# }

SESSION_STORE: dict[str, dict] = {}

# =========================
# Models
# =========================

class ChatRequest(BaseModel):
    message: str
    session_id: str | None = None
    user_name: str | None = None


class ChatResponse(BaseModel):
    reply: str


class UploadResponse(BaseModel):
    session_id: str
    status: str


# =========================
# Health Check
# =========================

@app.get("/health")
def health():
    return {"status": "ok"}


# =========================
# FILE UPLOAD (EPHEMERAL)
# =========================

@app.post("/upload", response_model=UploadResponse)
async def upload_document(file: UploadFile = File(...)):
    session_id = str(uuid.uuid4())
    logger.info(f"Upload started | session={session_id} | file={file.filename}")

    with timed("File read"):
        content = await file.read()

    with timed("Parsing"):
        raw_text = parse_file(file.filename, content)

    is_image = file.filename.lower().endswith(
        (".png", ".jpg", ".jpeg", ".bmp", ".tiff")
    )

    # =========================
    # 🖼 IMAGE: OCR-ONLY MODE
    # =========================
    if is_image:
        logger.info(
            f"Image detected | session={session_id} | ocr_len={len(raw_text.strip())}"
        )

        SESSION_STORE[session_id] = {
            "type": "image",
            "ocr_text": raw_text.strip(),
        }

        return UploadResponse(
            session_id=session_id,
            status="Image uploaded (OCR-only mode)",
        )

    # =========================
    # 📄 DOCUMENT PIPELINE
    # =========================
    if not raw_text.strip():
        logger.warning("No text extracted from document")
        raw_text = "No readable text found."

    text_len = len(raw_text)

    # 🔥 Adaptive chunk sizing
    if text_len < 3_000:
        chunk_size, overlap = 500, 50
    elif text_len < 15_000:
        chunk_size, overlap = 800, 100
    else:
        chunk_size, overlap = 1_200, 150

    logger.info(
        f"Chunking strategy | text_len={text_len} | "
        f"chunk_size={chunk_size} | overlap={overlap}"
    )

    with timed("Chunking"):
        chunks = chunk_text(
            raw_text,
            chunk_size=chunk_size,
            overlap=overlap,
        )

    with timed("Embedding"):
        embeddings = embed_chunks(chunks)

    SESSION_STORE[session_id] = {
        "type": "document",
        "chunks": chunks,
        "embeddings": embeddings,
    }

    logger.info(
        f"Document ingested | session={session_id} | chunks={len(chunks)}"
    )

    return UploadResponse(
        session_id=session_id,
        status="Document parsed and indexed",
    )


# =========================
# CHAT ENDPOINT (RAG-AWARE)
# =========================

@app.post("/chat", response_model=ChatResponse)
def chat(req: ChatRequest):
    logger.info(
        f"Chat request | session={req.session_id} | message_len={len(req.message)}"
    )

    store = SESSION_STORE.get(req.session_id)

    # =========================
    # 🖼 IMAGE CHAT
    # =========================
    if store and store.get("type") == "image":
        prompt = build_rag_prompt(
            user_message=req.message,
            context_chunks=[store.get("ocr_text", "")],
            user_name=req.user_name,
            mode="image",
            max_context_chars=2_000,
        )

    # =========================
    # 📄 DOCUMENT CHAT (RAG)
    # =========================
    elif store and store.get("type") == "document":
        with timed("Query embedding"):
            query_embedding = embed_chunks([req.message])[0]

        total_chunks = len(store["chunks"])
        top_k = min(5 + math.floor(total_chunks / 10), 8)

        logger.info(f"Retrieval | total_chunks={total_chunks} | top_k={top_k}")

        with timed("Retrieval"):
            retrieved_chunks = retrieve_relevant_chunks(
                query_embedding=query_embedding,
                chunk_embeddings=store["embeddings"],
                chunks=store["chunks"],
                top_k=top_k,
            )

        prompt = build_rag_prompt(
            user_message=req.message,
            context_chunks=retrieved_chunks,
            user_name=req.user_name,
            max_context_chars=6_000,
        )

    # =========================
    # 💬 NORMAL CHAT (NO FILE)
    # =========================
    else:
        prompt = build_rag_prompt(
            user_message=req.message,
            context_chunks=[],
            user_name=req.user_name,
        )

    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
    }

    try:
        with timed("LLM generation"):
            response = requests.post(
                OLLAMA_URL,
                json=payload,
                timeout=None,  # ⛔ No timeout for local LLM
            )
            response.raise_for_status()

        data = response.json()

        return ChatResponse(
            reply=data.get("response", "").strip()
        )

    except Exception:
        logger.exception("LLM generation failed")
        return ChatResponse(
            reply="⚠️ Cortex encountered an internal error while thinking."
        )


# =========================
# FILE PARSERS
# =========================

def parse_file(filename: str, content: bytes) -> str:
    filename = filename.lower()

    if filename.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")

    if filename.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if filename.endswith(".docx"):
        from docx import Document
        doc = Document(BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)

    if filename.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    if filename.endswith((".png", ".jpg", ".jpeg", ".bmp", ".tiff")):
        from PIL import Image
        import pytesseract

        pytesseract.pytesseract.tesseract_cmd = (
            r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        )

        img = Image.open(BytesIO(content)).convert("L")
        return pytesseract.image_to_string(img)

    return ""