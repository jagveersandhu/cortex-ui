from pptx import Presentation

def parse(file):
    prs = Presentation(file.file)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

    return "\n".join(text_runs)